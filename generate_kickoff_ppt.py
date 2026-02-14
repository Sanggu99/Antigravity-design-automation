from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define Colors
    TEAL = RGBColor(0, 128, 128)
    DARK_BLUE = RGBColor(10, 25, 50)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    ACCENT_ORANGE = RGBColor(255, 165, 0)

    # Helper function to set slide background
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper function to add Title Slide
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, DARK_BLUE)

        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(54)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_shape = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
        tf = sub_shape.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_ORANGE
        p.alignment = PP_ALIGN.CENTER
        
        # Decoration line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.333), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL
        shape.line.fill.background()

    # Helper function to add Content Slide
    def add_content_slide(title_text, content_items, image_guide, note=None):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, WHITE)

        # Title Bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = DARK_BLUE
        shape.line.width = Pt(1.5)
        
        # Title Text
        text_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(10), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.color.rgb = DARK_BLUE
        p.font.bold = True

        # Content Box (Left)
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = "\u2022 " + item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(14)
            p.space_before = Pt(6)

        # Image Placeholder (Right)
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.333), Inches(5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = LIGHT_GRAY
        placeholder.line.color.rgb = TEAL
        placeholder.line.width = Pt(1)
        
        # Image Text
        img_tf = placeholder.text_frame
        img_tf.word_wrap = True
        p = img_tf.paragraphs[0]
        p.text = "[IMAGE GUIDE]\n" + image_guide
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        
        # Note/Footer
        if note:
             footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11), Inches(0.5))
             tf = footer.text_frame
             p = tf.paragraphs[0]
             p.text = "NOTE: " + note
             p.font.size = Pt(12)
             p.font.italic = True
             p.font.color.rgb = TEAL

    # Slide 1: Title
    add_title_slide("서서울문화플라자 건립 설계공모", "Kick-off Meeting | 2026. 02.")

    # Slide 2: Table of Contents
    add_content_slide("목차", [
        "1. Project Overview (공모 개요)",
        "2. Context & Site Analysis (대지 분석)",
        "3. Regulations & Guidelines (법규 및 지침)",
        "4. Design Strategy (설계 전략 및 주안점)",
        "5. Schedule & Team (일정 및 수행계획)"
    ], "우측 절반에 대상지의 위성지도 클로즈업 또는 대지 경계선이 강조된 심플한 다이어그램 배치")

    # Slide 3: Project Understanding (Slogan)
    add_content_slide("프로젝트 이해 (Slogan)", [
        "\"서남권의 새로운 복합 문화 거점\"",
        "도서관, 스포츠, 키즈카페의 유기적 결합",
        "공공성 + 개방성 + 지속가능성 실현"
    ], "세 가지 주요 프로그램(책장, 수영장 레인, 아이들 놀이공간)이 겹쳐지거나 연결되는 추상적인 콜라주 또는 스케치 이미지", "팀의 강력한 의지를 보여주는 슬라이드")

    # Chapter 1
    add_content_slide("Chapter 1. 공모 개요", [
        "명칭: 서서울문화플라자 건립 설계공모",
        "방식: 일반설계공모",
        "주최: 서울특별시 (문화시설과)",
        "위치: 강서구 내발산동 743번지"
    ], "서울시 지도 내에서 강서구의 위치, 그리고 강서구 내에서 내발산동 대상지의 위치를 줌인하는 인포그래픽", "정확한 위치 인지")

    add_content_slide("사업 규모 및 예산", [
        "대지면적: 10,017.6㎡",
        "연면적: 11,000㎡ (±3% 조정 가능)",
        "층수: 지상 5층 이하",
        "예정공사비: 47,502백만원 (VAT 포함)",
        "설계비: 3,036백만원 (VAT 포함)"
    ], "면적 비교를 위한 인포그래픽 바 차트. 예산 규모를 보여주는 도식", "공사비 준수 중요")

    add_content_slide("주요 프로그램 구성", [
        "서울도서관 분관 (서남권 거점)",
        "생활체육센터 (수영장, 헬스장 등)",
        "서울형 키즈카페 (실내외 놀이공간)",
        "3가지 기능의 유기적 결합 요구"
    ], "이소메트릭 다이어그램으로 3개 매스가 하나로 합쳐지는 개념도. 각 매스에 라벨링 (Library, Sports, Kids)", "복합화가 키포인트")
    
    add_content_slide("공모 일정", [
        "참가등록: ~ 2026.03.06",
        "작품접수: 2026.04.10 (17:00 마감)",
        "1차 심사: 2026.04.22",
        "2차 심사(발표): 2026.04.29",
        "결과발표: 2026.05.06"
    ], "전체 일정을 가로로 긴 타임라인 형태로 표현. 현재 시점(Kick-off)과 마감일까지 남은 기간을 시각적으로 강조 (D-Day 카운트다운)", "마감 시간 엄수 (17:00)")

    add_content_slide("심사위원 및 평가 주안점", [
        "심사위원: 이기옥, 곽상준, 박미숙, 박현진, 이영기, 임진수, 정영신",
        "평가 배점:",
        "  - 목적 적합성(30)",
        "  - 복합화 전략(20)",
        "  - 도시맥락(20)",
        "  - 운영계획(20)",
        "  - 지속가능성(10)"
    ], "심사위원 명단 리스트와 배점 비율을 보여주는 도넛 차트 또는 레이더 차트", "심사위원 성향 분석 필요")

    # Chapter 2
    add_content_slide("Chapter 2. 대지 위치 및 현황", [
        "위치: 서울특별시 강서구 내발산동 743번지",
        "지구: 발산택지개발지구 지구단위계획구역 내 문화체육시설 부지",
        "현황: 주변 대규모 아파트 단지와 초중고교 밀집 지역"
    ], "대상지 반경 500m, 1km 위성지도. 주변 아파트 단지와 학교, 공원 위치를 아이콘으로 표시", "주거밀집지역의 커뮤니티 거점 역할 중요")

    add_content_slide("도시계획 및 지구단위계획", [
        "지역지구: 제2종일반주거지역, 중요시설물보호지구(공항)",
        "건폐율: 60% 이하",
        "용적률: 200% 이하",
        "높이: 5층 이하 (중요시설물보호지구 높이제한 별도 확인 필)"
    ], "지구단위계획 결정도면 스캔본. 대상지 부분 하이라이트. 허용 용도와 불허 용도를 보여주는 표", "공항 인근 고도제한 체크 필수")

    add_content_slide("접근성 및 동선 분석", [
        "차량 접근: 남측 소로 2-2호선(3m)? 주변 도로 현황 확인 필요",
        "보행 접근: 인근 주거단지에서의 도보 접근 경로",
        "대중교통: 지하철역 및 버스정류장 연계성 검토"
    ], "대지 주변 도로망 분석 다이어그램. 차량 진입 가능위치와 주 보행 동선 흐름을 화살표로 표시. 등하교 시간대 학생 동선 강조", "차량/보행 분리 계획 필수")

    add_content_slide("주변 맥락 (Context)", [
        "공공연계: 발산1동주민센터 등 인접 공공시설과의 연계 고려",
        "프라이버시: 주거지역 프라이버시 침해 방지 조치 필요",
        "일조권: 남측 일조권 확보에 유리한지 주변 건물 높이 등 분석"
    ], "대상지 주변 파노라마 사진 또는 3D 매스 모델링 캡쳐. 주변 건물의 높이와 배치를 보여주는 분석도", "민원 발생 소지 사전 차단")

    add_content_slide("자연환경 분석 (일조 및 바람)", [
        "향(Orientation): 남향 위주 배치 유리",
        "바람길(Wind Path): 자연환기를 고려한 통풍 계획",
        "지형(Level): 대지 레벨차 확인 (평탄지 여부 검토)"
    ], "태양 궤적도(Sun Path)를 오버랩한 배치 분석도. 바람장미(Wind Rose) 다이어그램", "친환경 디자인의 기초")

    add_content_slide("법규 제한사항 정리", [
        "건폐율: 60% (법정) → 계획시 여유공간 확보 가능",
        "용적률: 200% → 계획 연면적(11,000㎡) 달성 용이",
        "주차: 법정 주차대수 이상 여유주차대수 최대 확보 요구"
    ], "법적 최대치와 본 공모 요구치를 비교하는 막대 그래프 (건폐/용적률 여유분 시각화)", "넓은 야외공간 및 옥외주차 활용 가능성")

    add_content_slide("SWOT 분석", [
        "강점(Strength): 충분한 대지면적, 확실한 배후수요",
        "약점(Weakness): 공항 고도제한, 소음 가능성",
        "기회(Opportunity): 서남권 대표 랜드마크 가능성",
        "위협(Threat): 예산 내 시공성 확보, 물가 상승"
    ], "대지 사진 위에 강점/약점/기회/위협 요소를 텍스트와 화살표로 매핑한 종합 분석 다이어그램", "전략 도출의 근거")

    # Chapter 3
    add_content_slide("Chapter 3. 설계 기본 방향", [
        "\"서울도서관 분관으로서의 위상 확립\"",
        "\"3가지 시설(도서관+체육+키즈)의 유기적 융합\"",
        "\"열린(개방) 도서관 지향\""
    ], "개방형 도서관의 레퍼런스 이미지 (예: 헬싱키 오디 도서관 로비 등). 내부와 외부가 소통하는 투명한 파사드 스케치", "벽 없는 도서관 트렌드 반영")

    add_content_slide("배치 계획 주안점", [
        "유니버설 디자인: 모든 시민에게 차별 없는 이용 보장",
        "외부공간 활성화: 마당과 내부 프로그램의 유기적 연결",
        "동선 분리: 시설별 주차장/진입 동선의 합리적 분리 및 로비 통합"
    ], "대지 내 매스 배치 대안 스터디 다이어그램 (Alt 1, Alt 2, Alt 3). 외부공간(마당)을 중심으로 시설이 둘러싸는 배치 개념도", "외부공간 활성화 중요")

    add_content_slide("층별/시설별 조닝 계획", [
        "수직/수평 배치: 도서관/키즈카페/생활체육센터의 최적 배치",
        "연계: 공용공간(로비, 아트리움)을 통한 자연스러운 동선 유도",
        "분리: 소음 발생 시설(체육, 키즈)과 정숙 시설(도서관) 간 버퍼존 계획"
    ], "단면 조닝 다이어그램(Section Zoning). 층별로 색상을 달리하여 프로그램 배치 표현. 소음 차단 버퍼존 표시", "소음 관리가 핵심 기술 사항")

    add_content_slide("도서관 세부 지침", [
        "규모: 장서량 5만권 이상 (보존서고 12만권)",
        "구성: 일반자료실, 어린이/유아열람실, 미디어열람존, 북카페",
        "공간특성: 개방형 평면 구성, 서울엄마아빠VIP존 반영"
    ], "북큐레이션 공간, 계단식 열람 공간 등 최신 도서관 트렌드 무드보드", "가족 친화적 공간 조성")

    add_content_slide("생활체육센터 세부 지침", [
        "수영장: 25m x 5레인 (폭 2.3m↑), 수심 1.2~1.4m",
        "체육시설: 헬스장, GX룸, 다목적체육관(배드민턴 등)",
        "부대시설: 샤워실/탈의실 남녀 구분 및 충분한 규모 확보"
    ], "수영장 단면 상세 스케치 (지하 또는 저층부 배치 시 채광 유입 방안), 체육관의 높은 층고를 보여주는 스케치", "수영장 누수/결로 방지 및 구조 안정성")

    add_content_slide("키즈카페 세부 지침", [
        "기본모델: 서울형 키즈카페 가이드라인 적용",
        "연계: 실내 놀이공간 + 야외 놀이마당 연계 필수",
        "기능: 돌봄서비스 제공 공간, 다양한 놀이 환경 구축"
    ], "실내와 실외가 폴딩도어 등으로 연결되어 확장되는 놀이공간 개념도. 다채로운 색감의 놀이시설 레퍼런스", "안전 최우선 설계")

    add_content_slide("공용 및 지원시설 계획", [
        "코어: 통합로비, 북카페, 수유실 등 커뮤니티 중심 공간",
        "관리: 관리사무실, 통합방재센터의 효율적 배치",
        "주차: 옥내주차장 1,500㎡ 이상 (법정대수+@ 최대 확보)"
    ], "1층 평면 조닝 다이어그램. 통합 로비에서 각 시설로 분기되는 동선 체계(Wayfinding) 시각화", "로비는 커뮤니티의 중심")

    add_content_slide("친환경 및 에너지 계획", [
        "에너지 목표: 제로에너지빌딩(ZEB) 등급 목표 (지침 확인)",
        "인증: 녹색건축 인증 등 관련 인증 기준 준수",
        "신재생: 태양광 등 신재생에너지 적극 도입 및 디자인 요소화"
    ], "건물 외피(Facade)나 지붕에 적용된 태양광 패널 디자인 예시. 자연채광 및 자연환기 시뮬레이션 개념도", "지속가능한 친환경 디자인")

    add_content_slide("유니버설 디자인 (BF)", [
        "대상: 장애인, 노인, 임산부 등 모든 이용자",
        "계획: 단차 없는 진입, 충분한 통로 폭 확보",
        "인증: 장애물 없는 생활환경(BF) 인증 기준 준수 필"
    ], "무장애 동선 다이어그램. 휠체어 이동 시뮬레이션 라인", "레벨차 극복 방안 (경사로/EV)")

    add_content_slide("조경 및 외부공간 계획", [
        "녹지: 대지 내 충분한 조경면적 확보 및 옥상조경 활용",
        "개방: 주변 지역주민에게 개방된 휴식 공간 조성",
        "컨셉: 야외 도서관(Book Garden) 등 특화 공간 제안"
    ], "조경 식재 계획 개념도. 사계절 변화를 보여주는 외부공간 이미지. 야외 도서관(Book Garden) 컨셉 스케치", "도심 속 숲 같은 힐링 공간")
    
    # Chapter 4
    add_content_slide("Chapter 4. 설계 전략 (Concept)", [
        "Key Concept: [Culture Cloud / Forest Library / Interwoven Scape]",
        "전략 1: 경계 없는 융합 (Borderless Fusion)",
        "전략 2: 입체적 커뮤니티 (3D Community)",
        "전략 3: 자연을 품은 공간 (Eco-Friendly Void)"
    ], "프로젝트의 핵심 개념을 한 장으로 보여주는 강렬한 메인 컨셉 스케치 또는 다이어그램", "차별화된 디자인 포인트")

    add_content_slide("예상 이슈 및 해결방안", [
        "Issue 1: 이종 프로그램(정적 vs 동적)의 공존 → 해결: 층별 조닝/음향 시뮬레이션",
        "Issue 2: 공사비 상승 우려 → 해결: 합리적 구조 및 모듈화 계획",
        "Issue 3: 주차 공간 부족 → 해결: 효율적 주차 레이아웃 및 램프 계획"
    ], "이슈별 솔루션 매칭 테이블 또는 아이콘", "현실적 문제 해결 능력 어필")

    add_content_slide("수행 일정 및 인력 계획", [
        "일정: 디자인 스터디 → 중간평가 → 최종마감 등 마일스톤 관리",
        "인력: PM, 디자인, 도면, CG, 모형 등 분야별 전문 인력 투입",
        "협업: 구조, 기계, 전기, 조경 등 협력사와의 긴밀한 코디네이션"
    ], "간트 차트(Gantt Chart) 형태의 상세 작업 일정표. 조직도(Organization Chart)", "팀워크와 실행력 강조")

    add_content_slide("질의응답 (Q&A)", [
        "지침서 관련 모호한 점 토의",
        "발주처 요구사항 재확인",
        "향후 추진 방향 논의"
    ], "질문 아이콘 또는 심플한 배경", "자유 토론")

    # Closing Slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BLUE)
    
    text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.333), Inches(2))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다.\nThank You"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    prs.save('Kickoff_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
